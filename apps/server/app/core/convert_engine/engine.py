from typing import List, Optional
import os
import asyncio
from concurrent.futures import ThreadPoolExecutor
from loguru import logger

executor = ThreadPoolExecutor(max_workers=4)


class ConvertEngine:
    def __init__(self):
        pass
    
    async def pdf_to_word(
        self,
        input_path: str,
        output_path: str,
        quality: str = "high"
    ) -> str:
        def _convert():
            try:
                from pdf2docx import Converter
                
                cv = Converter(input_path)
                cv.convert(output_path)
                cv.close()
                
                return output_path
            except ImportError:
                logger.warning("pdf2docx not installed")
                raise RuntimeError("PDF to Word conversion requires pdf2docx package")
        
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(executor, _convert)
    
    async def pdf_to_excel(self, input_path: str, output_path: str) -> str:
        def _convert():
            try:
                import pdfplumber
                import pandas as pd
                
                all_tables = []
                with pdfplumber.open(input_path) as pdf:
                    for page in pdf.pages:
                        tables = page.extract_tables()
                        for table in tables:
                            df = pd.DataFrame(table[1:], columns=table[0])
                            all_tables.append(df)
                
                if all_tables:
                    with pd.ExcelWriter(output_path) as writer:
                        for i, df in enumerate(all_tables):
                            df.to_excel(writer, sheet_name=f'Table_{i+1}', index=False)
                
                return output_path
            except ImportError:
                logger.warning("pdfplumber or pandas not installed")
                raise RuntimeError("PDF to Excel conversion requires pdfplumber and pandas packages")
        
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(executor, _convert)
    
    async def pdf_to_ppt(self, input_path: str, output_path: str) -> str:
        def _convert():
            try:
                from pptx import Presentation
                from pptx.util import Inches
                import fitz
                
                prs = Presentation()
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)
                
                doc = fitz.open(input_path)
                
                for page in doc:
                    pix = page.get_pixmap(dpi=150)
                    img_data = pix.tobytes("png")
                    
                    import tempfile
                    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                        tmp.write(img_data)
                        tmp_path = tmp.name
                    
                    blank_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(blank_layout)
                    
                    slide.shapes.add_picture(
                        tmp_path,
                        Inches(0.5),
                        Inches(0.5),
                        width=Inches(12.333)
                    )
                    
                    os.unlink(tmp_path)
                
                doc.close()
                prs.save(output_path)
                
                return output_path
            except ImportError:
                logger.warning("python-pptx not installed")
                raise RuntimeError("PDF to PPT conversion requires python-pptx package")
        
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(executor, _convert)
    
    async def pdf_to_image(
        self,
        input_path: str,
        output_dir: str,
        image_format: str = "png",
        dpi: int = 150
    ) -> List[str]:
        def _convert():
            import fitz
            
            doc = fitz.open(input_path)
            output_files = []
            
            for i, page in enumerate(doc):
                pix = page.get_pixmap(dpi=dpi)
                output_file = os.path.join(output_dir, f"page_{i+1}.{image_format}")
                pix.save(output_file)
                output_files.append(output_file)
            
            doc.close()
            return output_files
        
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(executor, _convert)
    
    async def pdf_to_html(self, input_path: str, output_path: str) -> str:
        def _convert():
            import fitz
            
            doc = fitz.open(input_path)
            
            html_content = """<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>PDF Export</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 20px; }
        .page { margin-bottom: 20px; border-bottom: 1px solid #ccc; padding-bottom: 20px; }
    </style>
</head>
<body>
"""
            
            for i, page in enumerate(doc):
                text = page.get_text()
                html_content += f'<div class="page"><h3>Page {i+1}</h3><pre>{text}</pre></div>\n'
            
            html_content += "</body></html>"
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            doc.close()
            return output_path
        
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(executor, _convert)
    
    async def word_to_pdf(self, input_path: str, output_path: str) -> str:
        def _convert():
            try:
                from docx import Document
                from reportlab.lib.pagesizes import letter
                from reportlab.pdfgen import canvas
                from reportlab.lib.units import inch
                
                doc = Document(input_path)
                
                c = canvas.Canvas(output_path, pagesize=letter)
                width, height = letter
                
                y = height - inch
                
                for para in doc.paragraphs:
                    text = para.text
                    if text.strip():
                        if y < inch:
                            c.showPage()
                            y = height - inch
                        
                        c.drawString(inch, y, text[:100])
                        y -= 14
                
                c.save()
                return output_path
            except ImportError:
                logger.warning("python-docx or reportlab not installed")
                raise RuntimeError("Word to PDF conversion requires python-docx and reportlab packages")
        
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(executor, _convert)
    
    async def excel_to_pdf(self, input_path: str, output_path: str) -> str:
        def _convert():
            try:
                import openpyxl
                from reportlab.lib.pagesizes import letter
                from reportlab.pdfgen import canvas
                from reportlab.lib.units import inch
                
                wb = openpyxl.load_workbook(input_path)
                
                c = canvas.Canvas(output_path, pagesize=letter)
                width, height = letter
                
                y = height - inch
                
                for sheet in wb.worksheets:
                    c.drawString(inch, y, f"Sheet: {sheet.title}")
                    y -= 20
                    
                    for row in sheet.iter_rows(values_only=True):
                        if y < inch:
                            c.showPage()
                            y = height - inch
                        
                        row_text = ' | '.join(str(cell or '') for cell in row)
                        c.drawString(inch, y, row_text[:80])
                        y -= 14
                
                c.save()
                return output_path
            except ImportError:
                logger.warning("openpyxl or reportlab not installed")
                raise RuntimeError("Excel to PDF conversion requires openpyxl and reportlab packages")
        
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(executor, _convert)
    
    async def images_to_pdf(self, input_paths: List[str], output_path: str) -> str:
        def _convert():
            import fitz
            
            doc = fitz.open()
            
            for img_path in input_paths:
                img_doc = fitz.open(img_path)
                rect = img_doc[0].rect
                pdf_bytes = img_doc.convert_to_pdf()
                img_doc.close()
                
                img_pdf = fitz.open("pdf", pdf_bytes)
                page = doc.new_page(width=rect.width, height=rect.height)
                page.show_pdf_page(rect, img_pdf, 0)
                img_pdf.close()
            
            doc.save(output_path)
            doc.close()
            return output_path
        
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(executor, _convert)
    
    async def html_to_pdf(self, input_path: str, output_path: str) -> str:
        def _convert():
            try:
                from weasyprint import HTML
                
                HTML(input_path).write_pdf(output_path)
                return output_path
            except ImportError:
                logger.warning("weasyprint not installed")
                raise RuntimeError("HTML to PDF conversion requires weasyprint package")
        
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(executor, _convert)
